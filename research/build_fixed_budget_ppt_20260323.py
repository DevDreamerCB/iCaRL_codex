from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


OUT = Path("/data1/bochen/cbcontinual/iCaRL_codex/research/fixed_budget_progression_20260323_zh.pptx")


def add_bullets(slide, title, bullets, subtitle=None):
    title_box = slide.shapes.title
    title_box.text = title
    if subtitle is not None:
        tx = slide.placeholders[1].text_frame
        tx.clear()
        p = tx.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        for item in bullets:
            para = tx.add_paragraph()
            para.text = item
            para.level = 0
            para.font.size = Pt(22)
    else:
        left = Inches(0.8)
        top = Inches(1.4)
        width = Inches(11.5)
        height = Inches(5.4)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for item in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(22)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    (
        "固定预算标准下的实验递进",
        [
            "主标准：stage_epochs = 10,12,12",
            "screen：1 seed；full：3 seeds",
            "目标：在统一预算下比较旧任务保持与新任务学习的平衡",
            "当前最可信主线：s2ace00 + new_only ptblend",
        ],
        "跨被试类增量学习：AB -> BC -> CD",
    ),
    (
        "为什么要改标准",
        [
            "最近一批 full 回到了默认 30/30/30，混入了训练时长效应。",
            "重新检查后，10,12,12 口径下当前主线的 task3 明显更高。",
            "连续学习论文更适合用固定预算比较，而不是各方法都长训到收敛。",
        ],
        None,
    ),
    (
        "最核心的场景判断",
        [
            "task2 的主要瓶颈是 subjects 4/5/6 的弱域适应，而不是单纯 forgetting。",
            "joint(4,5,6) 上界说明：A=Left Hand 在新域并不是学不会。",
            "因此真正该修的是：stage2 当前样本里，缺失旧类 A 被如何对待。",
        ],
        None,
    ),
    (
        "主线递进 1",
        [
            "baseline_confirm：85.80 / 59.85 / 37.23",
            "no_contrastive_confirm：85.80 / 60.78 / 40.24",
            "no_contrastive_adapter16_s2_confirm：85.34 / 61.52 / 41.34",
            "说明：去掉 contrastive + 轻量 adapter 是最早的有效起点。",
        ],
        None,
    ),
    (
        "主线递进 2",
        [
            "lwf015_replaymix2_adapter16_mem36_oldweight25_confirm：85.73 / 58.49 / 42.72",
            "lwf015_normnme_adapter16_mem36_oldweight25_confirm：86.11 / 59.46 / 44.85",
            "lwf015_normnme_adapter16_mem36_oldweight30_stage101212_confirm：84.18 / 60.54 / 47.79",
            "说明：oldweight + normNME + 合理 stage budget 让非 hybrid 主线成型。",
        ],
        None,
    ),
    (
        "主线递进 3",
        [
            "非 hybrid 成熟线：84.18 / 60.86 / 48.29",
            "旧 fixed-budget best：84.18 / 60.75 / 49.36",
            "当前 fixed-budget 主线：84.18 / 60.70 / 49.40",
            "关键变化：从 stage3 specialist 逐渐转向更通用的 stage2 机制。",
        ],
        None,
    ),
    (
        "当前最佳方法",
        [
            "运行名：s2ace00_ptblendnew02_stage101212_confirm",
            "结果：84.18 / 60.70 / 49.40",
            "核心组成：s2ace00 + normNME + hybrid + new_only ptblend + adapter16 + affine + LwF + stage3 FD",
            "优点：更通用、更可解释，比继续堆 stage3 技巧更适合论文叙事。",
        ],
        None,
    ),
    (
        "为什么 s2ace00 有价值",
        [
            "s2ace00 是 stage2 asymmetric BCE：缺失旧类 A 不再在当前 B/C 样本上承受错误负压力。",
            "它直接针对 task2 的核心缺口，而不是再去修 stage3 的局部 trade-off。",
            "和 balanced finetune、stage2 replay trick 相比，它更简洁，也更有文献支撑。",
        ],
        None,
    ),
    (
        "近期尝试但未成立的方法",
        [
            "s2curce family：0.1 short 很强，但 0.1/0.05 的 full 都崩。",
            "stage2 balanced finetune：方向有启发，但不如主线稳。",
            "stage2-only feature distill 前移：没有让方法更通用。",
            "单纯增加 stage2 训练预算到 14：没有提升 task2，反而破坏平衡。",
        ],
        None,
    ),
    (
        "简化主线的结论",
        [
            "noaffine full：84.18 / 61.57 / 49.58，说明 affine 有帮助但不是主驱动。",
            "adapter8 full：85.96 / 61.39 / 49.11，可作更轻版本，但不是最佳。",
            "adapter12 short 明显不行。",
            "因此当前更像：保留 adapter16，是否保留 affine 取决于论文取舍。",
        ],
        None,
    ),
    (
        "关于 epochs 的最新结论",
        [
            "同主线在 10,12,12 下：84.72 / 60.96 / 50.54（short）",
            "同主线在 30-epoch 口径下：近期结果反而混入了训练时长效应。",
            "现在推荐：主表统一用 10,12,12；30/30/30 只作为补充收敛性检查。",
        ],
        None,
    ),
    (
        "后续更值得做的方向",
        [
            "继续围绕 stage2 的缺失旧类处理，而不是继续叠 stage3 patch。",
            "优先尝试：更温和、更稳的 stage2 old-class treatment。",
            "评价时固定看：task2、S4/S5/S6、A=Left Hand on 4/5/6、task3。",
        ],
        None,
    ),
]


for idx, (title, bullets, subtitle) in enumerate(slides):
    if idx == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
    else:
        slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_bullets(slide, title, bullets, subtitle)

prs.save(OUT)
print(str(OUT))

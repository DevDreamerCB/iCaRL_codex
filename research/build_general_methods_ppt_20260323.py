from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt


OUT = Path("/data1/bochen/cbcontinual/iCaRL_codex/research/general_methods_chain_20260323_zh.pptx")


def add_slide(prs, title, bullets, subtitle=None):
    layout = prs.slide_layouts[1] if subtitle else prs.slide_layouts[5]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title

    if subtitle:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        for item in bullets:
            para = tf.add_paragraph()
            para.text = item
            para.level = 0
            para.font.size = Pt(22)
        return

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(11.6), Inches(5.6))
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    (
        "通用方法链复验",
        [
            "纠正上一版错误底座：这次从 mem36 + replaymix2 出发",
            "统一目标：只看全 task 通用组件，不用 stage-specific trick",
            "方法链仍按固定预算逐步叠加，但起点更合理",
        ],
        "更合理底座：no contrastive + memory=36 + replaymix2",
    ),
    (
        "为什么要重跑",
        [
            "上一版通用链错误地从 mem24 + 无 replaymix 出发。",
            "这会把很多本来依赖更合理 replay 底座的方法误判为无效。",
            "因此必须从更接近真实强基线的底座重新做一遍链式验证。",
        ],
        None,
    ),
    (
        "方法链步骤",
        [
            "ugen00：更合理底座（no contrastive + mem36 + replaymix2）",
            "ugen01：+ task adapter dim=16，从 stage1 开始",
            "ugen02：+ stage_epochs = 10,12,12",
            "ugen03：+ oldweight = 3.125",
            "ugen04：+ normalized NME",
            "ugen05：+ LwF(0.175, T=1.5)",
            "ugen06：+ task affine，从 stage1 开始",
            "ugen07：+ hybrid NME+logits，统一早期开启",
        ],
        None,
    ),
    (
        "结果总览",
        [
            "ugen00: 84.95 / 57.56 / 42.71",
            "ugen01: 85.42 / 57.48 / 42.59",
            "ugen02: 85.42 / 60.19 / 41.47",
            "ugen03: 85.42 / 52.39 / 39.81",
            "ugen04: 85.42 / 56.64 / 44.95",
            "ugen05: 85.42 / 55.86 / 46.22",
            "ugen06: 84.26 / 55.94 / 45.06",
            "ugen07: 84.26 / 55.17 / 45.99",
        ],
        None,
    ),
    (
        "底座本身说明什么",
        [
            "仅仅把底座改成 mem36 + replaymix2，就比上一版错误链更接近真实场景。",
            "ugen00 已经到 84.95 / 57.56 / 42.71。",
            "这证明之前的大差距主要来自起点过弱，而不是方法链本身全都失效。",
        ],
        None,
    ),
    (
        "adapter 全开并不成立",
        [
            "ugen01 只有 task1 略升，task2/task3 没有真实改善。",
            "说明 adapter16 从 stage1 全开，不是一个自然成立的全 task 通用组件。",
            "它更像条件性有效的方法，需要更具体的上下文。",
        ],
        None,
    ),
    (
        "10,12,12 的作用",
        [
            "ugen02 把 task2 从 57.48 拉到 60.19。",
            "说明更合理的阶段训练预算仍然是有效的通用设置。",
            "但它本身不是万能修正，task3 仍然偏低。",
        ],
        None,
    ),
    (
        "为什么 oldweight 不适合写成通用组件",
        [
            "ugen03 变成 85.42 / 52.39 / 39.81。",
            "这说明 oldweight=3.125 不能机械地全局加入通用链。",
            "它只有在更具体、更条件性的主线里才成立。",
        ],
        None,
    ),
    (
        "normNME 与 LwF 的地位",
        [
            "ugen04: normNME -> 85.42 / 56.64 / 44.95",
            "ugen05: +LwF -> 85.42 / 55.86 / 46.22",
            "说明这两项在更合理底座上仍然是最像“通用 continual 组件”的方法。",
        ],
        None,
    ),
    (
        "为什么 affine 和 hybrid 也没成立",
        [
            "ugen06 affineall: 84.26 / 55.94 / 45.06",
            "ugen07 hybridall: 84.26 / 55.17 / 45.99",
            "把它们从 stage1/早期统一启用，并不会自然形成增益。",
            "所以它们更适合被归类为条件性组件，而不是通用主线组件。",
        ],
        None,
    ),
    (
        "当前通用链结论",
        [
            "更合理底座：mem36 + replaymix2",
            "成立的通用点：10,12,12、normalized NME、LwF",
            "不成立的通用点：adapter 全开、全局 oldweight、affine 全开、hybrid 早期开启",
        ],
        None,
    ),
    (
        "和更强主线的关系",
        [
            "更强 fixed-budget 主线仍然是 s2ace00_ptblendnew02_stage101212_confirm。",
            "通用链回答的是“哪些组件能写成普适方法”。",
            "增强主线回答的是“针对当前场景瓶颈，哪些阶段性方法最有效”。",
            "这两套材料在论文里应该分层叙述，而不是混在一起。",
        ],
        None,
    ),
]

for title, bullets, subtitle in slides:
    add_slide(prs, title, bullets, subtitle)

prs.save(OUT)
print(str(OUT))
